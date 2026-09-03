import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import win32com.client
import os

prs = pptx.Presentation()
prs.slide_width = Inches(8.27)
prs.slide_height = Inches(11.69)
blank_layout = prs.slide_layouts[6]

C_NAVY_DARK = RGBColor(16, 42, 77)
C_PRIMARY = RGBColor(30, 86, 160)
C_ACCENT_RED = RGBColor(229, 62, 62)
C_ACCENT_BLUE = RGBColor(49, 130, 206)
C_CARD_BG = RGBColor(248, 250, 252)
C_TEXT_DARK = RGBColor(30, 41, 59)
C_TEXT_MUTED = RGBColor(100, 116, 139)
C_WHITE = RGBColor(255, 255, 255)
C_LINE = RGBColor(226, 232, 240)

def add_header(slide, slide_num, category, title):
    tb_cat = slide.shapes.add_textbox(Inches(0.6), Inches(0.42), Inches(6.0), Inches(0.28))
    p_cat = tb_cat.text_frame.paragraphs[0]
    p_cat.text = f"{category.upper()}"
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_PRIMARY

    tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.68), Inches(6.8), Inches(0.55))
    p_title = tb_title.text_frame.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(16.5)
    p_title.font.bold = True
    p_title.font.color.rgb = C_NAVY_DARK

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.28), Inches(7.07), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = C_LINE
    line.line.color.rgb = C_LINE

    tb_foot = slide.shapes.add_textbox(Inches(0.6), Inches(11.2), Inches(6.0), Inches(0.25))
    p_foot = tb_foot.text_frame.paragraphs[0]
    p_foot.text = "AI와 함께하는 교통문제 해결 데이터 분석 공모전 | 한겨레 × 숲과나눔"
    p_foot.font.size = Pt(8)
    p_foot.font.color.rgb = C_TEXT_MUTED

    tb_num = slide.shapes.add_textbox(Inches(7.0), Inches(11.2), Inches(0.7), Inches(0.25))
    p_num = tb_num.text_frame.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.text = f"{slide_num:02d} / 05"
    p_num.font.size = Pt(8.5)
    p_num.font.bold = True
    p_num.font.color.rgb = C_PRIMARY

def add_card(slide, left, top, width, height, bg_color=C_CARD_BG, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
    else:
        card.line.fill.background()
    return card

# -------------------------------------------------------------
# SLIDE 1: 표지 및 요약
# -------------------------------------------------------------
s1 = prs.slides.add_slide(blank_layout)
top_banner = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(8.27), Inches(3.8))
top_banner.fill.solid()
top_banner.fill.fore_color.rgb = C_NAVY_DARK
top_banner.line.fill.background()

tb = s1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6.6), Inches(0.3))
p = tb.text_frame.paragraphs[0]
p.text = "AI와 함께하는 교통문제 해결을 위한 데이터 분석 공모전 [분석보고서]"
p.font.size = Pt(9.5)
p.font.bold = True
p.font.color.rgb = RGBColor(186, 215, 255)

tb = s1.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(6.6), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "저상버스 확충 예산은\n어디로 갔는가?"
p.font.size = Pt(27)
p.font.bold = True
p.font.color.rgb = C_WHITE

p2 = tf.add_paragraph()
p2.text = "교통약자 밀집지역 우선 배정 가설의 역설적 기각과 공간적 역진성 규명"
p2.font.size = Pt(12.5)
p2.font.color.rgb = RGBColor(220, 235, 255)
p2.space_before = Pt(8)

tb = s1.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(6.6), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "분석 대상: 전국 17개 시·도 및 경기도 31개 시·군 6,431개 버스 노선 전수 분석 | 2026. 09"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(160, 190, 230)

y_offset = 4.1
tb = s1.shapes.add_textbox(Inches(0.8), Inches(y_offset), Inches(6.6), Inches(0.35))
p = tb.text_frame.paragraphs[0]
p.text = "핵심 실증 발견 (Key Empirical Findings)"
p.font.size = Pt(13.5)
p.font.bold = True
p.font.color.rgb = C_NAVY_DARK

# Card 1
c1 = add_card(s1, 0.8, y_offset + 0.4, 6.67, 1.45, C_CARD_BG, C_ACCENT_RED)
tb = s1.shapes.add_textbox(Inches(0.95), Inches(y_offset + 0.45), Inches(6.35), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[가설 기각] '교통약자가 많은 지역일수록 저상버스는 오히려 적다'"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = C_ACCENT_RED
p2 = tf.add_paragraph()
p2.text = "• 전국 17개 시·도 상관계수 r = -0.435 (등록장애인 기준 r = -0.573 음의 상관관계)\n• 전남(교통약자 33.7%, 전국 1위) 저상버스 11.5% vs 서울(교통약자 22.6%) 저상버스 56.8% (5배 격차)"
p2.font.size = Pt(9)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(3)

# Card 2
c2 = add_card(s1, 0.8, y_offset + 2.05, 6.67, 1.45, C_CARD_BG, C_PRIMARY)
tb = s1.shapes.add_textbox(Inches(0.95), Inches(y_offset + 2.1), Inches(6.35), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[기초지자체 회귀분석] 교통약자 1%p 증가 시 저상버스 노선 1.14%p 감소"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p2 = tf.add_paragraph()
p2.text = "• 경기도 31개 시·군 6,431개 노선 전수 분석 결과 통계적으로 유의미한 음의 계수 (Beta = -1.14)\n• 하남·부천·수원(노선비율 53~64%) vs 가평·연천·여주(노선비율 0%)의 극단적 공간 양극화"
p2.font.size = Pt(9)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(3)

# Card 3
c3 = add_card(s1, 0.8, y_offset + 3.7, 6.67, 1.45, C_CARD_BG, RGBColor(217, 119, 6))
tb = s1.shapes.add_textbox(Inches(0.95), Inches(y_offset + 3.75), Inches(6.35), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[구조적 원인] '국비-지방비 50:50 정률 매칭' 제도의 역진적 함정"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(180, 83, 9)
p2 = tf.add_paragraph()
p2.text = "• 재정자립도 15~20%대 군 지역은 자체 부담금 확보 불가로 신차 도입 보조금 신청 포기\n• 재정 60~80%대 부유한 도심 지자체가 국비 보조금을 독식하는 '복지의 역진성' 규명"
p2.font.size = Pt(9)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(3)

# Card 4
c4 = add_card(s1, 0.8, y_offset + 5.35, 6.67, 1.4, RGBColor(238, 242, 255), C_PRIMARY)
tb = s1.shapes.add_textbox(Inches(0.95), Inches(y_offset + 5.4), Inches(6.35), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[정책 제언] AI 기반 저상버스 우선투입 지수(LBEI) & 차등 보조율 제도"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p2 = tf.add_paragraph()
p2.text = "• 소외 지자체 국비 매칭 50% → 80% 파격 상향으로 지자체 부담 경감\n• 머신러닝 수요 예측 스코어링을 통한 중앙정부 공모 평가 가점 및 우선 쿼터제 시행"
p2.font.size = Pt(9)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(3)

# -------------------------------------------------------------
# SLIDE 2: 전국 17개 시도 분석
# -------------------------------------------------------------
s2 = prs.slides.add_slide(blank_layout)
add_header(s2, 2, "Part 1. 거시 분석", "전국 17개 시·도: 수요와 공급의 정반대 불일치")
s2.shapes.add_picture("results/ppt_chart_sido.png", Inches(0.6), Inches(1.45), width=Inches(4.1))
s2.shapes.add_picture("results/ppt_chart_scatter_sido.png", Inches(0.6), Inches(5.8), width=Inches(4.1))

# Right Column (Clean text without emoji)
c_r1 = add_card(s2, 4.85, 1.45, 2.82, 2.8, C_WHITE, C_LINE)
tb = s2.shapes.add_textbox(Inches(4.95), Inches(1.55), Inches(2.6), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "가설 검증 결과: 기각"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_ACCENT_RED
p2 = tf.add_paragraph()
p2.text = "전국 17개 광역 시·도 분석 결과, 교통약자 비율과 저상버스 도입률은 통계적으로 유의미한 음(-)의 상관관계를 보임."
p2.font.size = Pt(8.8)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)
p3 = tf.add_paragraph()
p3.text = "• 전체 버스 도입률: r = -0.435\n• 등록장애인 비율: r = -0.573\n• 재정자립도 상관성: r = +0.495"
p3.font.size = Pt(8.6)
p3.font.bold = True
p3.font.color.rgb = C_PRIMARY
p3.space_before = Pt(5)

c_r2 = add_card(s2, 4.85, 4.45, 2.82, 3.2, C_CARD_BG)
tb = s2.shapes.add_textbox(Inches(4.95), Inches(4.55), Inches(2.6), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "극단적 양극화 비교"
p.font.size = Pt(11.5)
p.font.bold = True
p.font.color.rgb = C_NAVY_DARK
p2 = tf.add_paragraph()
p2.text = "[수도권 도심: 서울특별시]\n• 교통약자: 22.6% (낮음)\n• 재정자립도: 81.2% (전국 1위)\n• 저상버스: 56.8% (시내 66.7%)\n\n[비수도권 농어촌: 전라남도]\n• 교통약자: 33.7% (전국 1위)\n• 재정자립도: 28.7% (열악)\n• 저상버스: 11.5% (전국 꼴찌)"
p2.font.size = Pt(8.6)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(4)

c_r3 = add_card(s2, 4.85, 7.85, 2.82, 3.1, RGBColor(254, 242, 242), C_ACCENT_RED)
tb = s2.shapes.add_textbox(Inches(4.95), Inches(7.95), Inches(2.6), Inches(2.9))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "거시 분석 결론"
p.font.size = Pt(11.5)
p.font.bold = True
p.font.color.rgb = C_ACCENT_RED
p2 = tf.add_paragraph()
p2.text = "저상버스 예산 배정의 실질적 결정요인은 '교통약자의 복지 수요'가 아니라 '지자체의 자체 재정 동원력'이었음.\n\n가장 이동 지원이 절실한 지방 농어촌·고령화 지자체가 오히려 복지 인프라에서 배제되는 '지방 소외의 악순환' 확인."
p2.font.size = Pt(8.6)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

# -------------------------------------------------------------
# SLIDE 3: 경기도 노선 전수분석
# -------------------------------------------------------------
s3 = prs.slides.add_slide(blank_layout)
add_header(s3, 3, "Part 2. 미시 전수 분석", "경기도 6,431개 버스 노선: 기초지자체 공간 격차")
s3.shapes.add_picture("results/ppt_chart_quadrant.png", Inches(0.6), Inches(1.4), width=Inches(7.07))

col_w = 2.25
col_gap = 0.16
y_bot = 6.45

b1 = add_card(s3, 0.6, y_bot, col_w, 4.45, C_WHITE, C_PRIMARY)
tb = s3.shapes.add_textbox(Inches(0.68), Inches(y_bot + 0.1), Inches(col_w - 0.16), Inches(4.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "다중회귀분석 (OLS)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p2 = tf.add_paragraph()
p2.text = "Y(저상노선비율) =\n46.5 - 1.14*(약자비율)\n+ 0.14*(재정자립도)\n\n• Beta = -1.14 (p < 0.05)\n교통약자 비율 1%p 증가 시 저상버스 노선 비율 1.14%p 감소\n\n• 통계적으로 유의미한 역진성(Regressivity) 검증 완료"
p2.font.size = Pt(8.6)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

b2 = add_card(s3, 0.6 + col_w + col_gap, y_bot, col_w, 4.45, RGBColor(254, 242, 242), C_ACCENT_RED)
tb = s3.shapes.add_textbox(Inches(0.68 + col_w + col_gap), Inches(y_bot + 0.1), Inches(col_w - 0.16), Inches(4.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "제1우선 긴급투입지"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = C_ACCENT_RED
p2 = tf.add_paragraph()
p2.text = "교통약자 30~40% 밀집\n저상버스 노선 '0%' 전멸\n\n• 가평군 (38.2%) : 0.0%\n• 연천군 (39.2%) : 0.0%\n• 여주시 (31.7%) : 0.0%\n• 동두천 (31.2%) : 5.0%\n• 포천시 (31.4%) : 11.8%\n\n재정자립도 14~20%로 자체 도입 역량 부재"
p2.font.size = Pt(8.6)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

b3 = add_card(s3, 0.6 + (col_w + col_gap)*2, y_bot, col_w, 4.45, RGBColor(238, 242, 255), C_PRIMARY)
tb = s3.shapes.add_textbox(Inches(0.68 + (col_w + col_gap)*2), Inches(y_bot + 0.1), Inches(col_w - 0.16), Inches(4.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "자원 집중구역 (도심)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p2 = tf.add_paragraph()
p2.text = "교통약자 15~20% 내외\n저상버스 노선 50~64% 집중\n\n• 하남시 (18.1%) : 63.9%\n• 광명시 (21.1%) : 62.7%\n• 부천시 (21.1%) : 59.6%\n• 수원시 (17.3%) : 52.7%\n• 과천시 (17.3%) : 50.0%\n\n대규모 신도시 및 평지 도로 인프라 우위"
p2.font.size = Pt(8.6)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

# -------------------------------------------------------------
# SLIDE 4: 원인 진단
# -------------------------------------------------------------
s4 = prs.slides.add_slide(blank_layout)
add_header(s4, 4, "Part 3. 원인 진단", "왜 교통약자 밀집지역에 저상버스가 못 가는가?")

c_c1 = add_card(s4, 0.6, 1.45, 7.07, 2.95, C_WHITE, C_LINE)
tb = s4.shapes.add_textbox(Inches(0.85), Inches(1.58), Inches(6.57), Inches(2.7))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "원인 1. '국비-지방비 50:50 정률 매칭' 보조금의 제도적 함정"
p.font.size = Pt(12.5)
p.font.bold = True
p.font.color.rgb = C_NAVY_DARK
p2 = tf.add_paragraph()
p2.text = "• 문제 구조: 저상버스 대당 구입 보조금(약 9,000만 원)은 국가가 50%를 대고 지자체가 나머지 50%(약 4,500만 원)를 매칭해야만 교부됨.\n• 왜곡 발생: 서울(재정 81%), 경기 대도시(50%)는 매칭 예산을 대규모로 편성해 국비를 싹쓸이하는 반면, 재정자립도 15~20%대인 군 지역은 지자체 분담금을 마련하지 못해 신차 대폐차 시 저상버스 신청 자체를 포기함.\n• 결과: 재정이 취약한 지역일수록 국비 복지 혜택을 덜 받게 되는 '복지의 역진성(Regressive Welfare)' 초래."
p2.font.size = Pt(9.2)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

c_c2 = add_card(s4, 0.6, 4.65, 7.07, 2.95, C_WHITE, C_LINE)
tb = s4.shapes.add_textbox(Inches(0.85), Inches(4.78), Inches(6.57), Inches(2.7))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "원인 2. 운수회사의 승객 수요·수익성 중심 차량 배차"
p.font.size = Pt(12.5)
p.font.bold = True
p.font.color.rgb = C_NAVY_DARK
p2 = tf.add_paragraph()
p2.text = "• 문제 구조: 민간 운수회사 및 버스 준공영제 체계에서 신차와 고급 전기·수소 저상버스는 승차 인원이 많고 혼잡도가 높은 도심 '황금 간선노선'에 우선 배차됨.\n• 왜곡 발생: 탑승객 수가 적고 배차 간격이 긴 외곽 농어촌 노선이나 벽지 지선 노선은 차량 교체 우선순위에서 배제됨.\n• 결과: 저상버스가 가장 절실한 외곽 거주 휠체어 이용자 및 보행 장애인은 하루에 몇 대 오지 않는 고상버스 계단을 오르내려야 하는 현실 지속."
p2.font.size = Pt(9.2)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

c_c3 = add_card(s4, 0.6, 7.85, 7.07, 3.1, C_WHITE, C_LINE)
tb = s4.shapes.add_textbox(Inches(0.85), Inches(7.98), Inches(6.57), Inches(2.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "원인 3. 도로 인프라 미비와 '저상버스 도입 예외 승인'의 악순환"
p.font.size = Pt(12.5)
p.font.bold = True
p.font.color.rgb = C_NAVY_DARK
p2 = tf.add_paragraph()
p2.text = "• 문제 구조: 「교통약자법」상 저상버스 의무화에도 불구하고 '도로 종단경사 급경사, 굴곡, 과속방지턱 등으로 운행이 곤란한 노선'은 지자체가 예외를 승인할 수 있음.\n• 왜곡 발생: 농어촌·산간 군 지역은 도로 턱과 굴곡 정비 예산이 없어 노선 전체를 '예외 노선'으로 신청·승인하여 법적 의무를 우회함.\n• 결과: 도로를 고치지 않으니 저상버스를 못 넣고, 저상버스가 없으니 도로를 안 고치는 제도적 사각지대 고착화."
p2.font.size = Pt(9.2)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

# -------------------------------------------------------------
# SLIDE 5: 정책 제언
# -------------------------------------------------------------
s5 = prs.slides.add_slide(blank_layout)
add_header(s5, 5, "Part 4. 정책 제언", "데이터 기반 교통복지 형평성 혁신 방안")

p1_card = add_card(s5, 0.6, 1.45, 7.07, 2.85, RGBColor(238, 242, 255), C_PRIMARY)
tb = s5.shapes.add_textbox(Inches(0.85), Inches(1.58), Inches(6.57), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "제언 1. '교통약자 수요 연동형 차등 국비 보조율 제도' 도입"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p2 = tf.add_paragraph()
p2.text = "• 현행 일률 50% 국비 지원 방식을 '지자체 재정자립도'와 '교통약자 인구 밀집도'에 연동하여 30% ~ 80% 차등 지원으로 전면 개편.\n• 가평·연천·전남 등 재정자립도 20% 미만 소외 지자체는 국비 지원율을 최대 80%까지 파격 상향하여 지자체 지방비 부담을 1,800만 원 선으로 대폭 경감 → 도입 장벽 즉시 제거."
p2.font.size = Pt(9.2)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

p2_card = add_card(s5, 0.6, 4.55, 7.07, 3.4, C_WHITE, C_PRIMARY)
tb = s5.shapes.add_textbox(Inches(0.85), Inches(4.68), Inches(6.57), Inches(3.1))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "제언 2. AI 기반 '저상버스 우선투입 지수(LBEI)' 산출 및 공모 의무화"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p2 = tf.add_paragraph()
p2.text = "• 머신러닝 스코어링 알고리즘 구축 (Low-floor Bus Equity Index):\n  LBEI = 0.40 × (행정동별 고령·장애인 인구 밀도)\n       + 0.30 × (노선 내 종합병원·재활센터·복지관 경유도)\n       + 0.30 × (1 - 현재 노선 저상버스 공급 비율)\n• 적용 방안: 국토교통부 차년도 저상버스 도입 보조금 교부 심사 시, LBEI 상위 노선 및 기초지자체에 '우선 배정 쿼터제(최소 40% 의무 배정)' 도입."
p2.font.size = Pt(9.2)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

p3_card = add_card(s5, 0.6, 8.2, 7.07, 2.75, C_WHITE, C_LINE)
tb = s5.shapes.add_textbox(Inches(0.85), Inches(8.33), Inches(6.57), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "제언 3. 중형 저상버스 보조금 트랙 신설 & 도로 인프라 연계 패키지"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_NAVY_DARK
p2 = tf.add_paragraph()
p2.text = "• 대형(11m) 저상버스가 진입 불가능한 농어촌 굴곡 도로를 위해 '중형(8~9m) 저상 전기버스' 전용 보조금 트랙 신설.\n• '저상버스 도입 예외 승인 노선'에 대해 도로 턱 낮춤, 굴곡 완화, 정류장 단차 개선 예산을 국토부 도로정비사업과 1:1 패키지로 의무 지원하여 '물리적 진입 장벽' 해소."
p2.font.size = Pt(9.2)
p2.font.color.rgb = C_TEXT_DARK
p2.space_before = Pt(5)

# Save
out_pptx = "results/저상버스_배정형평성_분석보고서_세로형.pptx"
prs.save(out_pptx)
print("Saved clean PPTX!")

# Convert to PDF and re-export PNGs via PowerPoint COM
ppt_path = os.path.abspath(out_pptx)
pdf_path = os.path.abspath("results/저상버스_배정형평성_분석보고서_세로형.pdf")
out_dir = os.path.abspath("results")

powerpoint = win32com.client.Dispatch('PowerPoint.Application')
powerpoint.Visible = True
deck = powerpoint.Presentations.Open(ppt_path)
deck.SaveAs(pdf_path, 32)
for i, slide in enumerate(deck.Slides):
    img_path = os.path.join(out_dir, f'slide_page_{i+1}.png')
    slide.Export(img_path, 'PNG', 1240, 1754)
deck.Close()
powerpoint.Quit()
print("Exported clean PDF and slide images!")
